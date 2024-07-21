# List all shapes and details on specified slide
from pptx import Presentation


def list_slide_shapes(pptx_path, slide_index):
    # Iterate through all shapes on the slide and print details
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]

    for i, shape in enumerate(slide.shapes):
        print(f"Shape {i + 1}:")
        print(f"  Type: {shape.shape_type}")

        if shape.has_text_frame:
            print(f"  Text: {shape.text}")

        print(
            f"  Position: Left={shape.left}, Top={shape.top}, Width={shape.width}, Height={shape.height}"
        )
        print("-" * 30)


def main():
    # from pptx import Presentation

    # Path to PowerPoint presentation
    pptx_path = "pptx_tables.pptx"

    # Specify the slide index you want to analyze (adjust index as needed)
    slide_index = 0

    list_slide_shapes(pptx_path, slide_index)


if __name__ == "__main__":
    main()
